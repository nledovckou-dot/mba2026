#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Скрипт для конвертации HTML презентации в PPTX
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import re

def clean_text(text):
    """Очистка текста от HTML тегов и лишних пробелов"""
    if not text:
        return ""
    # Удаляем HTML теги
    text = re.sub(r'<[^>]+>', '', str(text))
    # Удаляем лишние пробелы
    text = ' '.join(text.split())
    return text.strip()

def extract_slide_content(slide_html):
    """Извлечение содержимого слайда из HTML"""
    soup = BeautifulSoup(str(slide_html), 'html.parser')
    
    # Извлекаем заголовок
    title_elem = soup.find(class_='slide-title')
    title = clean_text(title_elem.get_text()) if title_elem else ""
    
    # Извлекаем подзаголовок
    subtitle_elem = soup.find(class_='slide-subtitle')
    subtitle = clean_text(subtitle_elem.get_text()) if subtitle_elem else ""
    
    # Извлекаем все списки
    lists = []
    for ul in soup.find_all(['ul', 'ol']):
        items = []
        for li in ul.find_all('li'):
            items.append(clean_text(li.get_text()))
        if items:
            lists.append(items)
    
    # Извлекаем текст из блоков
    content_blocks = []
    for block in soup.find_all(['div', 'p', 'span']):
        classes = block.get('class', [])
        if any(c in ['key-point', 'highlight-box', 'stat-box', 'competitor-box', 'collaboration-item', 'utp-box'] for c in classes):
            text = clean_text(block.get_text())
            if text and len(text) > 10:  # Игнорируем очень короткие блоки
                content_blocks.append(text)
    
    return {
        'title': title,
        'subtitle': subtitle,
        'lists': lists,
        'content_blocks': content_blocks
    }

def create_pptx_from_html(html_file, output_file):
    """Создание PPTX из HTML файла"""
    
    # Читаем HTML
    with open(html_file, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Находим все слайды
    slides_html = soup.find_all('div', class_='slide')
    
    # Создаем презентацию
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Цвета ВШЭ
    hse_blue = RGBColor(26, 35, 126)  # #1a237e
    
    for slide_html in slides_html:
        content = extract_slide_content(slide_html)
        
        # Создаем слайд с пустым макетом
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой слайд
        
        # Добавляем заголовок
        if content['title']:
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.text = content['title']
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = hse_blue
            title_para.alignment = PP_ALIGN.LEFT
        
        # Добавляем подзаголовок
        if content['subtitle']:
            left = Inches(0.5)
            top = Inches(1.8)
            width = Inches(9)
            height = Inches(0.8)
            
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = content['subtitle']
            subtitle_para.font.size = Pt(32)
            subtitle_para.font.color.rgb = hse_blue
            subtitle_para.alignment = PP_ALIGN.LEFT
        
        # Добавляем списки
        current_top = Inches(2.8)
        for list_items in content['lists']:
            if current_top > Inches(6.5):
                break  # Не помещаемся на слайд
            
            left = Inches(0.7)
            width = Inches(9)
            height = Inches(min(0.3 * len(list_items), 4))
            
            list_box = slide.shapes.add_textbox(left, current_top, width, height)
            list_frame = list_box.text_frame
            list_frame.word_wrap = True
            
            for i, item in enumerate(list_items[:15]):  # Ограничиваем количество элементов
                if i == 0:
                    para = list_frame.paragraphs[0]
                else:
                    para = list_frame.add_paragraph()
                
                para.text = f"• {item}"
                para.font.size = Pt(14)
                para.space_after = Pt(6)
                para.level = 0
            
            current_top += height + Inches(0.2)
        
        # Добавляем текстовые блоки
        for block in content['content_blocks'][:3]:  # Ограничиваем количество блоков
            if current_top > Inches(6.5):
                break
            
            left = Inches(0.5)
            width = Inches(9)
            height = Inches(0.5)
            
            block_box = slide.shapes.add_textbox(left, current_top, width, height)
            block_frame = block_box.text_frame
            block_frame.word_wrap = True
            block_para = block_frame.paragraphs[0]
            block_para.text = block[:200]  # Ограничиваем длину
            block_para.font.size = Pt(12)
            block_para.space_after = Pt(8)
            
            current_top += Inches(0.7)
    
    # Сохраняем презентацию
    prs.save(output_file)
    print(f"✅ Презентация создана: {output_file}")

if __name__ == "__main__":
    html_file = "Анализ_рынка_презентация.html"
    output_file = "Анализ_рынка_презентация.pptx"
    
    try:
        create_pptx_from_html(html_file, output_file)
    except ImportError:
        print("❌ Необходимо установить библиотеки:")
        print("pip3 install python-pptx beautifulsoup4")
    except Exception as e:
        print(f"❌ Ошибка: {e}")

